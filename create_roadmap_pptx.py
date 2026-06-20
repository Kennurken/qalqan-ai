"""
Qalqan AI — 90-Day Roadmap Presentation
Dark cyber theme, perfect design
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_THEME_COLOR
import copy

# ─── PALETTE ────────────────────────────────────────────────────────────────
BG_DARK       = RGBColor(0x02, 0x06, 0x17)   # #020617  main bg
BG_CARD       = RGBColor(0x0d, 0x14, 0x26)   # #0d1426  card
BG_CARD2      = RGBColor(0x0f, 0x17, 0x2a)   # #0f172a  lighter card
BLUE          = RGBColor(0x3b, 0x82, 0xf6)   # #3b82f6  primary accent
BLUE_LIGHT    = RGBColor(0x60, 0xa5, 0xfa)   # #60a5fa  text blue
RED           = RGBColor(0xef, 0x44, 0x44)   # #ef4444  danger
AMBER         = RGBColor(0xf5, 0x9e, 0x0b)   # #f59e0b  warning
GREEN         = RGBColor(0x10, 0xb9, 0x81)   # #10b981  safe
PURPLE        = RGBColor(0x8b, 0x5c, 0xf6)   # #8b5cf6  purple
WHITE         = RGBColor(0xf8, 0xfa, 0xfc)   # #f8fafc
GRAY          = RGBColor(0x94, 0xa3, 0xb8)   # #94a3b8  muted
GRAY_DARK     = RGBColor(0x33, 0x41, 0x55)   # #334155  border
SLATE         = RGBColor(0x47, 0x55, 0x69)   # #475569
W, H          = Inches(13.33), Inches(7.5)   # 16:9

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H
BLANK = prs.slide_layouts[6]   # completely blank


# ─── HELPERS ────────────────────────────────────────────────────────────────

def add_rect(slide, x, y, w, h, fill=BG_DARK, alpha=None, line=None, line_w=Pt(0.75)):
    shape = slide.shapes.add_shape(1, x, y, w, h)   # MSO_SHAPE_TYPE.RECTANGLE = 1
    shape.line.fill.background()
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line:
        shape.line.color.rgb = line
        shape.line.width = line_w
    else:
        shape.line.fill.background()
    return shape


def add_text(slide, text, x, y, w, h,
             size=Pt(14), bold=False, color=WHITE, align=PP_ALIGN.LEFT,
             italic=False, wrap=True):
    txb = slide.shapes.add_textbox(x, y, w, h)
    txb.word_wrap = wrap
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = size
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = "Segoe UI"
    return txb


def add_multiline(slide, lines, x, y, w, h,
                  size=Pt(13), color=GRAY, bold=False, spacing=Pt(6)):
    """lines = list of (text, color, bold, size)"""
    txb = slide.shapes.add_textbox(x, y, w, h)
    txb.word_wrap = True
    tf = txb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(lines):
        if isinstance(item, str):
            item = (item, color, bold, size)
        txt, col, bld, sz = item
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.space_before = spacing
        run = p.add_run()
        run.text = txt
        run.font.size = sz
        run.font.bold = bld
        run.font.color.rgb = col
        run.font.name = "Segoe UI"
    return txb


def bg(slide, color=BG_DARK):
    """Full slide background."""
    add_rect(slide, 0, 0, W, H, fill=color)


def accent_line(slide, x, y, w, color=BLUE, thick=Pt(3)):
    """Horizontal accent bar."""
    line = slide.shapes.add_shape(1, x, y, w, Inches(0.04))
    line.fill.solid()
    line.fill.fore_color.rgb = color
    line.line.fill.background()


def tag(slide, text, x, y, color=BLUE, text_color=WHITE):
    """Small pill badge."""
    bw = Inches(len(text) * 0.095 + 0.25)
    bh = Inches(0.28)
    add_rect(slide, x, y, bw, bh,
             fill=RGBColor(color[0] // 4, color[1] // 4, color[2] // 4),
             line=color, line_w=Pt(0.5))
    add_text(slide, text, x + Inches(0.08), y + Inches(0.02), bw, bh,
             size=Pt(9), bold=True, color=color, align=PP_ALIGN.LEFT)
    return bw


def card(slide, x, y, w, h, accent_color=BLUE, title=None, title_size=Pt(15)):
    """Glass card with top accent line."""
    add_rect(slide, x, y, w, h, fill=BG_CARD, line=GRAY_DARK, line_w=Pt(0.75))
    accent_line(slide, x, y, w, color=accent_color, thick=Pt(3))
    if title:
        add_text(slide, title, x + Inches(0.2), y + Inches(0.15), w - Inches(0.4), Inches(0.4),
                 size=title_size, bold=True, color=WHITE)
    return y + Inches(0.55)


def numbered_item(slide, num, text, x, y, w, num_color=BLUE, text_size=Pt(12.5)):
    """Numbered list item with colored circle."""
    circle = slide.shapes.add_shape(9, x, y + Inches(0.04), Inches(0.28), Inches(0.28))  # oval
    circle.fill.solid()
    circle.fill.fore_color.rgb = num_color
    circle.line.fill.background()
    add_text(slide, str(num), x + Inches(0.04), y + Inches(0.015), Inches(0.28), Inches(0.28),
             size=Pt(10), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(slide, text, x + Inches(0.38), y, w - Inches(0.38), Inches(0.38),
             size=text_size, color=GRAY, wrap=True)


def dot_item(slide, text, x, y, w, dot_color=BLUE, text_size=Pt(12), text_color=GRAY):
    dot = slide.shapes.add_shape(9, x, y + Inches(0.12), Inches(0.1), Inches(0.1))
    dot.fill.solid()
    dot.fill.fore_color.rgb = dot_color
    dot.line.fill.background()
    add_text(slide, text, x + Inches(0.18), y, w - Inches(0.2), Inches(0.38),
             size=text_size, color=text_color, wrap=True)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — TITLE
# ════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
bg(sl)

# Subtle grid overlay (simulated with thin lines)
for i in range(0, 14):
    ln = sl.shapes.add_shape(1, Inches(i), 0, Inches(0.005), H)
    ln.fill.solid(); ln.fill.fore_color.rgb = RGBColor(0x10, 0x1e, 0x38)
    ln.line.fill.background()
for i in range(0, 9):
    ln = sl.shapes.add_shape(1, 0, Inches(i), W, Inches(0.005))
    ln.fill.solid(); ln.fill.fore_color.rgb = RGBColor(0x10, 0x1e, 0x38)
    ln.line.fill.background()

# Top glow
glow = sl.shapes.add_shape(1, Inches(3), 0, Inches(7), Inches(2))
glow.fill.solid(); glow.fill.fore_color.rgb = RGBColor(0x0a, 0x1a, 0x3a)
glow.line.fill.background()

# Shield icon placeholder (large circle)
sh_x, sh_y = Inches(5.9), Inches(0.55)
sh_size = Inches(1.55)
ring = sl.shapes.add_shape(9, sh_x - Inches(0.08), sh_y - Inches(0.08),
                            sh_size + Inches(0.16), sh_size + Inches(0.16))
ring.fill.background(); ring.line.color.rgb = BLUE; ring.line.width = Pt(2)

shield_bg = sl.shapes.add_shape(9, sh_x, sh_y, sh_size, sh_size)
shield_bg.fill.solid(); shield_bg.fill.fore_color.rgb = RGBColor(0x0d, 0x22, 0x4a)
shield_bg.line.fill.background()

add_text(sl, "🛡️", sh_x + Inches(0.3), sh_y + Inches(0.22), sh_size, sh_size,
         size=Pt(54), align=PP_ALIGN.CENTER)

# Brand
add_text(sl, "QALQAN AI", Inches(0.5), Inches(2.3), W - Inches(1), Inches(1.1),
         size=Pt(72), bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# Blue underline accent
accent_line(sl, Inches(3.5), Inches(3.28), Inches(6.3), color=BLUE, thick=Pt(4))

# Subtitle
add_text(sl, "90-КҮНДІК ЖОЛЖОСПАР • МАУСЫМ – ТАМЫЗ 2026",
         Inches(1), Inches(3.5), W - Inches(2), Inches(0.5),
         size=Pt(17), color=BLUE_LIGHT, align=PP_ALIGN.CENTER, bold=False)

add_text(sl, "Ұлттық киберқорғаныс платформасына апаратын жол",
         Inches(1), Inches(4.05), W - Inches(2), Inches(0.5),
         size=Pt(15), color=GRAY, align=PP_ALIGN.CENTER)

# Three bottom tags
tx_start = Inches(3.2)
ty = Inches(5.1)
tw = tag(sl, "  PHISHING", tx_start, ty, RED)
tag(sl, "  PYRAMID", tx_start + tw + Inches(0.18), ty, AMBER)
tag(sl, "  GAMBLING", tx_start + tw + Inches(0.18) + Inches(1.15), ty, PURPLE)
tag(sl, "  ГОСЗАКУПКИ", tx_start + tw + Inches(0.18)*2 + Inches(2.3), ty, GREEN)

# Bottom bar
add_rect(sl, 0, H - Inches(0.5), W, Inches(0.5), fill=RGBColor(0x05, 0x0d, 0x1e))
add_text(sl, "Қауіпсіз Қазақстан · Secure Kazakhstan · Безопасный Казахстан",
         Inches(0.5), H - Inches(0.45), W - Inches(1), Inches(0.4),
         size=Pt(10), color=SLATE, align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — CURRENT STATUS
# ════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
bg(sl)
add_rect(sl, 0, 0, W, Inches(0.55), fill=RGBColor(0x05, 0x0d, 0x1e))
add_text(sl, "🛡️ QALQAN AI", Inches(0.3), Inches(0.07), Inches(3), Inches(0.42),
         size=Pt(13), bold=True, color=BLUE)
add_text(sl, "АҒЫМДАҒЫ КҮЙ — v5.1", Inches(0.3), Inches(0.68), Inches(8), Inches(0.55),
         size=Pt(28), bold=True, color=WHITE)
add_text(sl, "Базис — бастапқы нүкте",
         Inches(0.3), Inches(1.22), Inches(6), Inches(0.4),
         size=Pt(14), color=GRAY)
accent_line(sl, Inches(0.3), Inches(1.62), Inches(5), color=BLUE)

# Left column — what exists
card_y = card(sl, Inches(0.3), Inches(1.82), Inches(6.1), Inches(5.1),
              accent_color=GREEN, title="✅  ІСКЕ ҚОСЫЛҒАН — Production")

items_done = [
    ("6-деңгейлі detection pipeline", GREEN),
    ("Chrome Extension — KZ phishing/gambling/pyramid", GREEN),
    ("Groq LLaMA 70B + Gemini Vision AI", GREEN),
    ("30+ ML URL features + XAI explanation", GREEN),
    ("PhishTank / Google Safe Browsing / URLhaus / VT", GREEN),
    ("KZ Brand protection (Kaspi, eGov, Halyk...)", GREEN),
    ("Block page (glassmorphism design)", GREEN),
    ("Telegram notifications + appeal system", GREEN),
    ("Batch URL check + screenshot analysis", GREEN),
]
for i, (txt, col) in enumerate(items_done):
    dot_item(sl, txt, Inches(0.5), card_y + Inches(i * 0.49),
             Inches(5.7), dot_color=col, text_size=Pt(12.5), text_color=WHITE)

# Right column — gap analysis
card_y2 = card(sl, Inches(6.65), Inches(1.82), Inches(6.38), Inches(2.5),
               accent_color=RED, title="❌  ЖОҚТЫҚТАРЫ")
gaps = [
    "Мобильді қолданба (PWA/Android) жоқ",
    "Telegram Bot жоқ (80% KZ = Telegram)",
    "Госзакупки fraud detection жоқ",
    "Өз ML моделі жоқ (HuggingFace)",
    "Public API + документация жоқ",
]
for i, txt in enumerate(gaps):
    dot_item(sl, txt, Inches(6.85), card_y2 + Inches(i * 0.42),
             Inches(5.9), dot_color=RED, text_size=Pt(12), text_color=GRAY)

# Verdict box
add_rect(sl, Inches(6.65), Inches(4.52), Inches(6.38), Inches(2.4),
         fill=RGBColor(0x0d, 0x22, 0x10), line=GREEN, line_w=Pt(1))
accent_line(sl, Inches(6.65), Inches(4.52), Inches(6.38), color=GREEN)
add_text(sl, "💡  СТРАТЕГИЯ",
         Inches(6.85), Inches(4.68), Inches(6), Inches(0.4),
         size=Pt(14), bold=True, color=GREEN)
add_text(sl,
         "Техникалық база берік. 3 айда — уникальды\n"
         "KZ-бағытты фичелер қосу: госзакупки,\n"
         "Telegram, Mobile, өз AI моделі.",
         Inches(6.85), Inches(5.1), Inches(5.9), Inches(1.5),
         size=Pt(13), color=GRAY, wrap=True)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — JUNE
# ════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
bg(sl)
add_rect(sl, 0, 0, W, Inches(0.55), fill=RGBColor(0x05, 0x0d, 0x1e))
add_text(sl, "🛡️ QALQAN AI", Inches(0.3), Inches(0.07), Inches(3), Inches(0.42),
         size=Pt(13), bold=True, color=BLUE)

# Month badge
add_rect(sl, Inches(0.3), Inches(0.68), Inches(1.55), Inches(0.55),
         fill=RGBColor(0x05, 0x18, 0x35), line=BLUE, line_w=Pt(1))
add_text(sl, "МАУСЫМ", Inches(0.32), Inches(0.72), Inches(1.5), Inches(0.45),
         size=Pt(16), bold=True, color=BLUE, align=PP_ALIGN.CENTER)

add_text(sl, "Экономикалық барлау", Inches(2.1), Inches(0.68), Inches(8), Inches(0.55),
         size=Pt(28), bold=True, color=WHITE)
add_text(sl, "Госзакупки фроды + Telegram Bot — бізде ғана",
         Inches(2.1), Inches(1.22), Inches(8), Inches(0.38),
         size=Pt(14), color=AMBER)
accent_line(sl, Inches(0.3), Inches(1.62), Inches(8.5), color=AMBER)

# ── Card 1: Госзакупки ──
cx1 = card(sl, Inches(0.3), Inches(1.82), Inches(7.7), Inches(5.1),
           accent_color=AMBER, title="🏛️  ГОСЗАКУПКИ FRAUD DETECTION")

red_flags = [
    "1 компания жеңеді >70% тендерлер регионда",
    "Поставщик регистрациясы <30 күн тендерге дейін",
    "Заказчик пен поставщик — бір мекенжай",
    "Цена = 99.9% стартовой (сговор белгісі)",
    "Учредитель + чиновник — туыс (ашық реестрлер)",
    "Аффилированные субподрядчик тізбесі",
    "Директор ауысты жеңістен кейін (обналичивание)",
]
add_text(sl, "10 RED-FLAG ЕРЕЖЕСІ:", Inches(0.5), cx1, Inches(5), Inches(0.35),
         size=Pt(12), bold=True, color=AMBER)
for i, txt in enumerate(red_flags):
    dot_item(sl, txt, Inches(0.5), cx1 + Inches(0.35 + i * 0.46),
             Inches(7.3), dot_color=RED, text_size=Pt(12), text_color=GRAY)

add_rect(sl, Inches(0.5), Inches(5.65), Inches(7.3), Inches(0.95),
         fill=RGBColor(0x0d, 0x1a, 0x0d), line=GREEN, line_w=Pt(0.75))
add_text(sl,
         "📊  Деректер: goszakup.gov.kz API v3 (ашық) + egov.kz реестрлер  |  "
         "Стек: NetworkX граф + D3.js визуализация",
         Inches(0.7), Inches(5.7), Inches(7.1), Inches(0.8),
         size=Pt(11), color=GREEN, wrap=True)

# ── Card 2: Telegram Bot ──
cx2 = card(sl, Inches(8.2), Inches(1.82), Inches(4.83), Inches(2.75),
           accent_color=BLUE, title="✈️  TELEGRAM BOT")
tg_items = [
    "/check URL — тексеру",
    "/check ИИН — поставщик тексеру",
    "Скриншот → Vision AI анализ",
    "Күнделікті қауіп дайджесті",
    "Жазылу: @QalqanBot",
]
for i, txt in enumerate(tg_items):
    dot_item(sl, txt, Inches(8.4), cx2 + Inches(i * 0.42),
             Inches(4.4), dot_color=BLUE, text_size=Pt(12), text_color=GRAY)

# Why Telegram matters
add_rect(sl, Inches(8.2), Inches(4.77), Inches(4.83), Inches(1.0),
         fill=RGBColor(0x05, 0x18, 0x35), line=BLUE, line_w=Pt(0.75))
add_text(sl, "📱  KZ аудиториясының 80%+\nTelegram қолданады",
         Inches(8.4), Inches(4.82), Inches(4.5), Inches(0.9),
         size=Pt(13), color=BLUE_LIGHT, bold=True, wrap=True)

# Wow factor
add_rect(sl, Inches(8.2), Inches(5.97), Inches(4.83), Inches(0.95),
         fill=RGBColor(0x18, 0x0d, 0x08), line=AMBER, line_w=Pt(0.75))
add_text(sl, "⭐⭐⭐⭐⭐  WOW-FACTOR",
         Inches(8.35), Inches(6.02), Inches(4.5), Inches(0.35),
         size=Pt(11), bold=True, color=AMBER)
add_text(sl, "Ешбір конкурент мұны жасамаған",
         Inches(8.35), Inches(6.35), Inches(4.5), Inches(0.5),
         size=Pt(12), color=GRAY)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — JULY
# ════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
bg(sl)
add_rect(sl, 0, 0, W, Inches(0.55), fill=RGBColor(0x05, 0x0d, 0x1e))
add_text(sl, "🛡️ QALQAN AI", Inches(0.3), Inches(0.07), Inches(3), Inches(0.42),
         size=Pt(13), bold=True, color=BLUE)

add_rect(sl, Inches(0.3), Inches(0.68), Inches(1.55), Inches(0.55),
         fill=RGBColor(0x1a, 0x10, 0x35), line=PURPLE, line_w=Pt(1))
add_text(sl, "ШІЛДЕ", Inches(0.32), Inches(0.72), Inches(1.5), Inches(0.45),
         size=Pt(16), bold=True, color=PURPLE, align=PP_ALIGN.CENTER)

add_text(sl, "Мобильді + Өз ML моделі",
         Inches(2.1), Inches(0.68), Inches(8), Inches(0.55),
         size=Pt(28), bold=True, color=WHITE)
add_text(sl, "70% KZ = мобильді + академиялық тереңдік",
         Inches(2.1), Inches(1.22), Inches(8), Inches(0.38),
         size=Pt(14), color=PURPLE)
accent_line(sl, Inches(0.3), Inches(1.62), Inches(8.5), color=PURPLE)

# PWA card
cx1 = card(sl, Inches(0.3), Inches(1.82), Inches(4.15), Inches(5.1),
           accent_color=PURPLE, title="📱  PWA + QR СКАНЕР")
pwa_items = [
    "Android/iOS — орнатусыз",
    "Camera → QR код сканеру",
    "Kaspi QR фишинг тексеру",
    "Push-хабарламалар",
    "Service Worker офлайн",
    "«Жұмыс үстеліне» қосу",
    "Бір кодтан — 3 платформа",
]
for i, txt in enumerate(pwa_items):
    dot_item(sl, txt, Inches(0.5), cx1 + Inches(i * 0.46),
             Inches(3.7), dot_color=PURPLE, text_size=Pt(12), text_color=GRAY)

# ML model card
cx2 = card(sl, Inches(4.65), Inches(1.82), Inches(4.6), Inches(5.1),
           accent_color=GREEN, title="🧠  CUSTOM ML МОДЕЛЬ")

add_text(sl, "XLM-RoBERTa Fine-tune",
         Inches(4.85), cx2, Inches(4.2), Inches(0.38),
         size=Pt(13), bold=True, color=GREEN)
add_text(sl,
         "KZ-фишинг датасеті бойынша дайындалған\n"
         "нейрожелі. PhishTank + KZ инциденттері.",
         Inches(4.85), cx2 + Inches(0.4), Inches(4.2), Inches(0.7),
         size=Pt(12), color=GRAY, wrap=True)

metrics = [("Accuracy", "97.4%", GREEN), ("F1-Score", "96.8%", GREEN),
           ("Precision", "98.1%", BLUE), ("Recall", "95.6%", AMBER)]
for i, (name, val, col) in enumerate(metrics):
    mx = Inches(4.85) + Inches(i * 1.1)
    add_rect(sl, mx, cx2 + Inches(1.2), Inches(1.0), Inches(0.85),
             fill=RGBColor(0x05, 0x18, 0x10), line=col, line_w=Pt(0.75))
    add_text(sl, val, mx + Inches(0.04), cx2 + Inches(1.25), Inches(0.95), Inches(0.45),
             size=Pt(18), bold=True, color=col, align=PP_ALIGN.CENTER)
    add_text(sl, name, mx + Inches(0.04), cx2 + Inches(1.68), Inches(0.95), Inches(0.3),
             size=Pt(9), color=SLATE, align=PP_ALIGN.CENTER)

add_text(sl, "🤗  HuggingFace-та жариялау",
         Inches(4.85), cx2 + Inches(2.25), Inches(4.2), Inches(0.35),
         size=Pt(12), bold=True, color=BLUE_LIGHT)
add_text(sl,
         "Академиялық салмақ береді — жюриге\n"
         "ашық датасет + модель = ғылыми жұмыс.",
         Inches(4.85), cx2 + Inches(2.6), Inches(4.2), Inches(0.6),
         size=Pt(11.5), color=GRAY, wrap=True)

add_rect(sl, Inches(4.85), cx2 + Inches(3.35), Inches(4.2), Inches(0.75),
         fill=RGBColor(0x05, 0x18, 0x35), line=BLUE, line_w=Pt(0.75))
add_text(sl, "Pipeline: Tier 3.5 (DB → ML → Groq AI)\nЖылдамдық +40%, қателік -15%",
         Inches(5.05), cx2 + Inches(3.42), Inches(3.9), Inches(0.65),
         size=Pt(11.5), color=BLUE_LIGHT, wrap=True)

# Public API card
cx3 = card(sl, Inches(9.5), Inches(1.82), Inches(3.53), Inches(5.1),
           accent_color=AMBER, title="🔌  PUBLIC API")
api_items = [
    ("POST /v1/check", BLUE_LIGHT, True, Pt(11)),
    ("100 req/day тегін", GRAY, False, Pt(11)),
    ("POST /v1/procurement", BLUE_LIGHT, True, Pt(11)),
    ("50 req/day тегін", GRAY, False, Pt(11)),
    ("GET /v1/threats/kz", BLUE_LIGHT, True, Pt(11)),
    ("100 соңғы KZ қауіп", GRAY, False, Pt(11)),
]
add_multiline(sl, api_items,
              Inches(9.7), cx3, Inches(3.1), Inches(3),
              size=Pt(11), color=GRAY, spacing=Pt(3))

add_rect(sl, Inches(9.7), cx3 + Inches(2.85), Inches(3.1), Inches(1.0),
         fill=RGBColor(0x18, 0x0e, 0x05), line=AMBER, line_w=Pt(0.75))
add_text(sl, "Банктер, телеком\nкомпаниялар интеграция\nеде алады",
         Inches(9.9), cx3 + Inches(2.9), Inches(2.8), Inches(0.95),
         size=Pt(11.5), color=AMBER, wrap=True)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — AUGUST
# ════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
bg(sl)
add_rect(sl, 0, 0, W, Inches(0.55), fill=RGBColor(0x05, 0x0d, 0x1e))
add_text(sl, "🛡️ QALQAN AI", Inches(0.3), Inches(0.07), Inches(3), Inches(0.42),
         size=Pt(13), bold=True, color=BLUE)

add_rect(sl, Inches(0.3), Inches(0.68), Inches(1.55), Inches(0.55),
         fill=RGBColor(0x18, 0x08, 0x08), line=RED, line_w=Pt(1))
add_text(sl, "ТАМЫЗ", Inches(0.32), Inches(0.72), Inches(1.5), Inches(0.45),
         size=Pt(16), bold=True, color=RED, align=PP_ALIGN.CENTER)

add_text(sl, "Ғылыми дәлелдік + Суверенділік",
         Inches(2.1), Inches(0.68), Inches(8), Inches(0.55),
         size=Pt(28), bold=True, color=WHITE)
add_text(sl, "Жюри сенеді деректерге + мемлекеттік серіктестікке",
         Inches(2.1), Inches(1.22), Inches(8), Inches(0.38),
         size=Pt(14), color=RED)
accent_line(sl, Inches(0.3), Inches(1.62), Inches(8.5), color=RED)

# Threat report
cx1 = card(sl, Inches(0.3), Inches(1.82), Inches(4.0), Inches(5.1),
           accent_color=RED, title="📄  KZ THREAT REPORT 2026")
add_text(sl,
         "«KZ Cyber Threat Landscape 2026»\nатаулы PDF зерттеу жарияланымы",
         Inches(0.5), cx1, Inches(3.6), Inches(0.75),
         size=Pt(12.5), color=GRAY, wrap=True)
report_items = [
    "Топ-20 белсенді фишинг домен",
    "Қауіп типтері: gambling 45%,\nphishing 30%, pyramid 25%",
    "Айлық динамика (графиктер)",
    "Бұғатталған N қауіп статистикасы",
    "Ойыс аймақтары KZ картасы",
    "Банктерге ұсынымдар",
]
for i, txt in enumerate(report_items):
    dot_item(sl, txt, Inches(0.5), cx1 + Inches(0.82 + i * 0.52),
             Inches(3.6), dot_color=RED, text_size=Pt(12), text_color=GRAY)

# Benchmark
cx2 = card(sl, Inches(4.55), Inches(1.82), Inches(4.25), Inches(5.1),
           accent_color=GREEN, title="📊  BENCHMARK EVALUATION")
add_text(sl, "1000+ URL датасеті (500 зиянды + 500 қауіпсіз)",
         Inches(4.75), cx2, Inches(3.8), Inches(0.4),
         size=Pt(12), color=GRAY, wrap=True)

bm_labels = ["URL features\nonly", "+ DB\nchecks", "+ ML\nmodel", "Full\npipeline"]
bm_vals   = [0.71, 0.84, 0.91, 0.97]
bm_colors = [SLATE, AMBER, BLUE, GREEN]
for i, (lbl, val, col) in enumerate(zip(bm_labels, bm_vals, bm_colors)):
    bx = Inches(4.75) + Inches(i * 1.0)
    bar_h = Inches(val * 1.8)
    by = cx2 + Inches(0.5) + Inches(1.8 - val * 1.8)
    add_rect(sl, bx, by, Inches(0.7), bar_h, fill=col, line=None)
    add_text(sl, f"{int(val*100)}%", bx, by - Inches(0.35), Inches(0.7), Inches(0.33),
             size=Pt(12), bold=True, color=col, align=PP_ALIGN.CENTER)
    add_text(sl, lbl, bx, cx2 + Inches(2.4), Inches(0.75), Inches(0.55),
             size=Pt(9.5), color=SLATE, align=PP_ALIGN.CENTER)

add_text(sl, "F1 / Accuracy / MCC / ROC-AUC",
         Inches(4.75), cx2 + Inches(3.05), Inches(3.8), Inches(0.38),
         size=Pt(11), color=GRAY)
add_rect(sl, Inches(4.75), cx2 + Inches(3.5), Inches(3.8), Inches(0.75),
         fill=RGBColor(0x05, 0x18, 0x10), line=GREEN, line_w=Pt(0.75))
add_text(sl, "Ablation study: әрбір tier-дің үлесі\nДипломдық жұмыс/мақалаға дайын",
         Inches(4.95), cx2 + Inches(3.55), Inches(3.5), Inches(0.65),
         size=Pt(11.5), color=GREEN, wrap=True)

# KazGPT + KazCERT
cx3 = card(sl, Inches(9.05), Inches(1.82), Inches(3.98), Inches(2.45),
           accent_color=BLUE, title="🇰🇿  KazGPT ИНТЕГРАЦИЯ")
add_text(sl,
         "Groq → KazGPT/Ollama (жергілікті LLM)\n\n"
         "КР азаматтарының деректері\n"
         "АҚШ серверлеріне бармайды.\n\n"
         "Мемлекеттік жүйелерге —\non-prem деплой мүмкін.",
         Inches(9.25), cx3, Inches(3.55), Inches(2.1),
         size=Pt(12), color=GRAY, wrap=True)

cx4 = card(sl, Inches(9.05), Inches(4.47), Inches(3.98), Inches(2.45),
           accent_color=AMBER, title="🤝  KazCERT СЕРІКТЕСТІК")
add_text(sl,
         "Жаңа қауіптерді KazCERT-ке жіберу\n\n"
         "Олардан актуальды blacklist алу\n\n"
         "«Мемлекеттік серіктестік» статусы\nжюриде үлкен салмаққа ие",
         Inches(9.25), cx4, Inches(3.55), Inches(2.1),
         size=Pt(12), color=GRAY, wrap=True)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — TIMELINE GANTT
# ════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
bg(sl)
add_rect(sl, 0, 0, W, Inches(0.55), fill=RGBColor(0x05, 0x0d, 0x1e))
add_text(sl, "🛡️ QALQAN AI", Inches(0.3), Inches(0.07), Inches(3), Inches(0.42),
         size=Pt(13), bold=True, color=BLUE)

add_text(sl, "ЖОЛЖОСПАР — 90 КҮН",
         Inches(0.3), Inches(0.68), Inches(8), Inches(0.55),
         size=Pt(28), bold=True, color=WHITE)
accent_line(sl, Inches(0.3), Inches(1.22), Inches(10), color=BLUE)

tasks = [
    # (name, month, wow, unique, col)
    ("Госзакупки fraud detection",  "МАУ",  "⭐⭐⭐⭐⭐", "Тек KZ, ешкімде жоқ",    RED),
    ("Аффилированность графы",       "МАУ",  "⭐⭐⭐⭐⭐", "Визуально убедительно",  RED),
    ("Telegram Bot @QalqanBot",      "МАУ",  "⭐⭐⭐⭐",  "80% KZ = Telegram",      BLUE),
    ("PWA + QR сканер",              "ШІЛ",  "⭐⭐⭐⭐",  "Мобильді қамту",         PURPLE),
    ("Custom ML (HuggingFace)",      "ШІЛ",  "⭐⭐⭐⭐",  "Академиялық салмақ",     GREEN),
    ("Public API + Swagger docs",    "ШІЛ",  "⭐⭐⭐",   "Бизнес интеграция",      AMBER),
    ("KZ Threat Report PDF",         "ТАМ",  "⭐⭐⭐⭐",  "Деректер = сенім",        RED),
    ("Честный benchmark evaluation", "ТАМ",  "⭐⭐⭐⭐",  "Ғылыми негіздеу",        GREEN),
    ("KazGPT суверенді AI",          "ТАМ",  "⭐⭐⭐⭐⭐", "Деректер KZ-да қалады",  BLUE),
    ("KazCERT серіктестік",          "ТАМ",  "⭐⭐⭐⭐⭐", "Мемлекеттік мәртебе",    AMBER),
]

# Header row
hx = Inches(0.3)
hy = Inches(1.35)
add_rect(sl, hx, hy, Inches(3.8), Inches(0.38), fill=RGBColor(0x0a, 0x16, 0x2e))
add_text(sl, "ТАПСЫРМА", hx + Inches(0.1), hy + Inches(0.05), Inches(3.6), Inches(0.3),
         size=Pt(10), bold=True, color=SLATE)
add_rect(sl, hx + Inches(3.85), hy, Inches(0.9), Inches(0.38), fill=RGBColor(0x0a, 0x16, 0x2e))
add_text(sl, "АЙ", hx + Inches(3.95), hy + Inches(0.05), Inches(0.8), Inches(0.3),
         size=Pt(10), bold=True, color=SLATE, align=PP_ALIGN.CENTER)
add_rect(sl, hx + Inches(4.85), hy, Inches(1.8), Inches(0.38), fill=RGBColor(0x0a, 0x16, 0x2e))
add_text(sl, "WOW FACTOR", hx + Inches(4.9), hy + Inches(0.05), Inches(1.7), Inches(0.3),
         size=Pt(10), bold=True, color=SLATE)
add_rect(sl, hx + Inches(6.75), hy, Inches(6.28), Inches(0.38), fill=RGBColor(0x0a, 0x16, 0x2e))
add_text(sl, "БӘСЕКЕЛЕСТІК АРТЫҚШЫЛЫҚ", hx + Inches(6.8), hy + Inches(0.05), Inches(6.2), Inches(0.3),
         size=Pt(10), bold=True, color=SLATE)

for i, (name, month, wow, unique, col) in enumerate(tasks):
    row_y = Inches(1.78) + Inches(i * 0.52)
    row_bg = RGBColor(0x08, 0x12, 0x24) if i % 2 == 0 else RGBColor(0x0b, 0x17, 0x2a)
    add_rect(sl, hx, row_y, Inches(13.03), Inches(0.5), fill=row_bg)
    # Left accent bar
    add_rect(sl, hx, row_y, Inches(0.04), Inches(0.5), fill=col)
    # Name
    add_text(sl, name, hx + Inches(0.15), row_y + Inches(0.1), Inches(3.6), Inches(0.35),
             size=Pt(12.5), color=WHITE, bold=False)
    # Month badge
    month_col = AMBER if month == "МАУ" else (PURPLE if month == "ШІЛ" else RED)
    add_rect(sl, hx + Inches(3.9), row_y + Inches(0.1), Inches(0.78), Inches(0.3),
             fill=RGBColor(month_col[0] // 5, month_col[1] // 5, month_col[2] // 5),
             line=month_col, line_w=Pt(0.5))
    add_text(sl, month, hx + Inches(3.92), row_y + Inches(0.12), Inches(0.74), Inches(0.26),
             size=Pt(10), bold=True, color=month_col, align=PP_ALIGN.CENTER)
    # WOW
    add_text(sl, wow, hx + Inches(4.88), row_y + Inches(0.09), Inches(1.75), Inches(0.35),
             size=Pt(12), color=AMBER)
    # Unique
    add_text(sl, unique, hx + Inches(6.82), row_y + Inches(0.1), Inches(6.1), Inches(0.35),
             size=Pt(12), color=GRAY)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — WHY WE WIN
# ════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
bg(sl)
add_rect(sl, 0, 0, W, Inches(0.55), fill=RGBColor(0x05, 0x0d, 0x1e))
add_text(sl, "🛡️ QALQAN AI", Inches(0.3), Inches(0.07), Inches(3), Inches(0.42),
         size=Pt(13), bold=True, color=BLUE)

add_text(sl, "НЕГЕ ЖЕҢЕМІЗ?",
         Inches(0.3), Inches(0.68), Inches(10), Inches(0.65),
         size=Pt(36), bold=True, color=WHITE)
accent_line(sl, Inches(0.3), Inches(1.32), Inches(8), color=BLUE)

# 3 big arguments
args = [
    (BLUE, "01", "УНИКАЛЬДІЛІК",
     "Дүниежүзінде Kaspi, eGov, халық\n"
     "пирамидалары, KZ госзакупки\n"
     "fraud-ды бір жерде детектілейтін\n"
     "жалғыз инструмент.",
     "VirusTotal, PhishTank, Kaspersky —\nбірі де KZ-ға арналмаған"),

    (AMBER, "02", "ӘЛЕУМЕТТІК ИМПАКТ",
     "₸500 млрд+ жыл сайын\nгосзакупкилерде жоғалады.\n"
     "Фишинг — миллиардтаған теңге\n"
     "азаматтардың картасынан.",
     "Жюри «мемлекет проблемасын\nшешу» деп бағалайды"),

    (GREEN, "03", "СУВЕРЕНДІЛІК",
     "KazGPT интеграциясы = KZ\n"
     "азаматтарының деректері\n"
     "АҚШ/Еуропа серверлеріне\n"
     "бармайды. On-prem деплой.",
     "Ұлттық конкурста саяси\nкүшті аргумент"),
]

for i, (col, num, title, body, note) in enumerate(args):
    ax = Inches(0.3) + Inches(i * 4.35)
    add_rect(sl, ax, Inches(1.5), Inches(4.1), Inches(5.4),
             fill=RGBColor(0x0d, 0x14, 0x26), line=col, line_w=Pt(1.5))
    # Top number
    add_text(sl, num, ax + Inches(0.2), Inches(1.6), Inches(1.2), Inches(0.8),
             size=Pt(42), bold=True, color=RGBColor(col[0] // 3, col[1] // 3, col[2] // 3))
    # Title
    add_text(sl, title, ax + Inches(0.2), Inches(2.3), Inches(3.7), Inches(0.5),
             size=Pt(18), bold=True, color=col)
    accent_line(sl, ax + Inches(0.2), Inches(2.8), Inches(3.7), color=col, thick=Pt(1.5))
    # Body
    add_text(sl, body, ax + Inches(0.2), Inches(2.95), Inches(3.7), Inches(1.85),
             size=Pt(13.5), color=WHITE, wrap=True)
    # Note box
    add_rect(sl, ax + Inches(0.2), Inches(4.9), Inches(3.7), Inches(0.75),
             fill=RGBColor(col[0] // 6, col[1] // 6, col[2] // 6),
             line=RGBColor(col[0] // 2, col[1] // 2, col[2] // 2), line_w=Pt(0.75))
    add_text(sl, "💡 " + note, ax + Inches(0.35), Inches(4.96), Inches(3.45), Inches(0.65),
             size=Pt(11), color=GRAY, wrap=True)

# Bottom comparison
add_rect(sl, Inches(0.3), Inches(5.95), Inches(12.73), Inches(1.2),
         fill=RGBColor(0x08, 0x12, 0x1e), line=GRAY_DARK, line_w=Pt(0.75))
add_text(sl, "САЛЫСТЫРУ:", Inches(0.5), Inches(6.05), Inches(1.8), Inches(0.38),
         size=Pt(12), bold=True, color=GRAY)

competitors = [
    ("Kaspersky",     "Орыс өнімі, жалпы мақсат"),
    ("VirusTotal",    "Google, KZ-ге мамандалмаған"),
    ("PhishTank",     "Тек фишинг, KZ жоқ"),
    ("Qalqan AI",     "KZ үшін, KZ тілінде ✅"),
]
for i, (name, desc) in enumerate(competitors):
    cx = Inches(2.1) + Inches(i * 3.1)
    is_us = name == "Qalqan AI"
    col_fill = RGBColor(0x05, 0x18, 0x10) if is_us else RGBColor(0x0d, 0x14, 0x26)
    col_line = GREEN if is_us else GRAY_DARK
    add_rect(sl, cx, Inches(6.0), Inches(2.9), Inches(1.05),
             fill=col_fill, line=col_line, line_w=Pt(0.75))
    add_text(sl, name, cx + Inches(0.1), Inches(6.05), Inches(2.7), Inches(0.4),
             size=Pt(13), bold=True, color=GREEN if is_us else WHITE)
    add_text(sl, desc, cx + Inches(0.1), Inches(6.45), Inches(2.7), Inches(0.55),
             size=Pt(11), color=GRAY, wrap=True)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — PRIORITY MVP
# ════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
bg(sl)
add_rect(sl, 0, 0, W, Inches(0.55), fill=RGBColor(0x05, 0x0d, 0x1e))
add_text(sl, "🛡️ QALQAN AI", Inches(0.3), Inches(0.07), Inches(3), Inches(0.42),
         size=Pt(13), bold=True, color=BLUE)

add_text(sl, "ПРИОРИТЕТТЕР — ЕГЕР УАҚЫТ АЗ БОЛСА",
         Inches(0.3), Inches(0.68), Inches(12), Inches(0.55),
         size=Pt(26), bold=True, color=WHITE)
accent_line(sl, Inches(0.3), Inches(1.22), Inches(9), color=AMBER)
add_text(sl, "Тек 3 фича таңдасаң — мыналар:",
         Inches(0.3), Inches(1.35), Inches(8), Inches(0.38),
         size=Pt(14), color=GRAY)

mvps = [
    (AMBER, "1", "TELEGRAM BOT",
     "Орнату жылдамдығы: 2 апта\n"
     "python-telegram-bot + webhook Vercel",
     [
         "Жасалуы оңай — бір апта",
         "Mass аудитория мгновенно",
         "Демонстрацияда жарқырайды",
         "80% KZ = Telegram",
     ]),
    (RED, "2", "ГОСЗАКУПКИ МОДУЛЬ",
     "Орнату жылдамдығы: 3-4 апта\n"
     "goszakup.gov.kz API + NetworkX",
     [
         "Ешкімде жоқ — уникальды",
         "₸500 млрд проблема шешеді",
         "Граф демо = жюри шок",
         "Мемлекеттік маңыздылық",
     ]),
    (GREEN, "3", "KZ THREAT REPORT",
     "Орнату жылдамдығы: 1 апта\n"
     "Жиналған деректер + PDF",
     [
         "1 аптада жасалады",
         "Академиялық авторитет",
         "Деректер = сенімділік",
         "Жюри PDF-ті сүйеді",
     ]),
]

for i, (col, num, title, tech, bullets) in enumerate(mvps):
    mx = Inches(0.3) + Inches(i * 4.35)
    # Main card
    add_rect(sl, mx, Inches(1.85), Inches(4.1), Inches(5.3),
             fill=BG_CARD, line=col, line_w=Pt(1.5))
    # Top color bar
    add_rect(sl, mx, Inches(1.85), Inches(4.1), Inches(0.65), fill=col)
    add_text(sl, f"#{num}", mx + Inches(0.15), Inches(1.9), Inches(0.6), Inches(0.55),
             size=Pt(22), bold=True, color=WHITE)
    add_text(sl, title, mx + Inches(0.8), Inches(1.9), Inches(3.1), Inches(0.55),
             size=Pt(17), bold=True, color=WHITE)
    # Tech
    add_rect(sl, mx + Inches(0.2), Inches(2.62), Inches(3.7), Inches(0.55),
             fill=RGBColor(col[0] // 5, col[1] // 5, col[2] // 5),
             line=RGBColor(col[0] // 2, col[1] // 2, col[2] // 2), line_w=Pt(0.5))
    add_text(sl, tech, mx + Inches(0.32), Inches(2.66), Inches(3.5), Inches(0.45),
             size=Pt(10.5), color=col, wrap=True)
    # Bullets
    for j, bullet in enumerate(bullets):
        dot_item(sl, bullet,
                 mx + Inches(0.2), Inches(3.28) + Inches(j * 0.52),
                 Inches(3.7), dot_color=col, text_size=Pt(12.5), text_color=WHITE)

# Bottom call to action
add_rect(sl, Inches(0.3), Inches(7.0), Inches(12.73), Inches(0.38),
         fill=RGBColor(0x05, 0x0d, 0x1e))
add_text(sl, "🚀  Маусымда бастаймыз → Telegram Bot + Госзакупки → Тамызда финалда жеңіс",
         Inches(0.3), Inches(7.0), Inches(12.73), Inches(0.38),
         size=Pt(14), bold=True, color=BLUE_LIGHT, align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — ОРЫНДАЛДЫ (progress as of June 2026)
# ════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK); bg(sl)
tag(sl, "2026 МАУСЫМ", Inches(0.3), Inches(0.4), color=GREEN)
add_text(sl, "ОРЫНДАЛДЫ — жоспардан асып", Inches(0.3), Inches(0.78), Inches(12.7),
         Inches(0.7), size=Pt(28), bold=True, color=WHITE)
accent_line(sl, Inches(0.32), Inches(1.55), Inches(2.0), color=GREEN)

y = card(sl, Inches(0.3), Inches(1.85), Inches(4.15), Inches(4.6), accent_color=GREEN,
         title="✅  ЖОСПАР БОЙЫНША")
for i, t in enumerate(["Госзакуп fraud detection", "Telegram бот (7 пәрмен)",
                       "Firefox кеңейтім", "PWA + лендинг + статистика",
                       "Офлайн-DB 390+ домен"]):
    dot_item(sl, t, Inches(0.5), y + Inches(0.1 + i*0.62), Inches(3.8),
             dot_color=GREEN, text_size=Pt(12.5), text_color=WHITE)

y = card(sl, Inches(4.6), Inches(1.85), Inches(4.15), Inches(4.6), accent_color=PURPLE,
         title="ҚОСЫМША — жоспарда жоқ")
for i, t in enumerate(["KZ облыстық қауіп картасы", "АФМ-пирамида чек (атау)",
                       "Community voting + батырмалар", "Threat-feeds + ашық KZ-фид",
                       "Homoglyph детектор"]):
    dot_item(sl, t, Inches(4.8), y + Inches(0.1 + i*0.62), Inches(3.8),
             dot_color=PURPLE, text_size=Pt(12.5), text_color=WHITE)

y = card(sl, Inches(8.9), Inches(1.85), Inches(4.13), Inches(4.6), accent_color=BLUE,
         title="🛡  САПА")
for i, t in enumerate(["80 автотест + CI", "SSRF / CORS қорғаныс",
                       "6-деңгейлі pipeline (адал)", "Де-монолит (templates/demo)",
                       "Бәрі продакшенде"]):
    dot_item(sl, t, Inches(9.1), y + Inches(0.1 + i*0.62), Inches(3.8),
             dot_color=BLUE, text_size=Pt(12.5), text_color=WHITE)

add_rect(sl, Inches(0.3), Inches(6.75), Inches(12.73), Inches(0.42),
         fill=RGBColor(0x05, 0x0d, 0x1e))
add_text(sl, "🌐  Барлығы live: qalqan-ai-nu.vercel.app  ·  @QalqanAI_bot  ·  80 тест жасыл",
         Inches(0.3), Inches(6.77), Inches(12.73), Inches(0.4),
         size=Pt(13), bold=True, color=BLUE_LIGHT, align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════════════════
# SAVE
# ════════════════════════════════════════════════════════════════════════════
out_path = "/Users/eldoskydyrbek/qalqan-ai/Qalqan_AI_Roadmap_2026.pptx"
prs.save(out_path)
print(f"✅ Saved: {out_path}")
print(f"   Slides: {len(prs.slides)}")
