"""
Qalqan AI — ДЭР республикалық байқауы: қорғау презентациясы (защитная дека).
Қазақ тілінде · dark cyber тема. Генерация: python3 create_defense_pptx.py
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ─── PALETTE (matches create_roadmap_pptx.py) ────────────────────────────────
BG_DARK   = RGBColor(0x02, 0x06, 0x17)
BG_CARD   = RGBColor(0x0d, 0x14, 0x26)
BG_CARD2  = RGBColor(0x0f, 0x17, 0x2a)
BLUE      = RGBColor(0x3b, 0x82, 0xf6)
CYAN      = RGBColor(0x00, 0xd4, 0xff)
BLUE_LIGHT= RGBColor(0x60, 0xa5, 0xfa)
RED       = RGBColor(0xef, 0x44, 0x44)
AMBER     = RGBColor(0xf5, 0x9e, 0x0b)
GREEN     = RGBColor(0x10, 0xb9, 0x81)
PURPLE    = RGBColor(0x8b, 0x5c, 0xf6)
WHITE     = RGBColor(0xf8, 0xfa, 0xfc)
GRAY      = RGBColor(0x94, 0xa3, 0xb8)
GRAY_DARK = RGBColor(0x33, 0x41, 0x55)
W, H      = Inches(13.33), Inches(7.5)

prs = Presentation()
prs.slide_width, prs.slide_height = W, H
BLANK = prs.slide_layouts[6]


# ─── HELPERS ─────────────────────────────────────────────────────────────────
def add_rect(slide, x, y, w, h, fill=BG_DARK, line=None, line_w=Pt(0.75)):
    s = slide.shapes.add_shape(1, x, y, w, h)
    s.line.fill.background()
    if fill:
        s.fill.solid(); s.fill.fore_color.rgb = fill
    else:
        s.fill.background()
    if line:
        s.line.color.rgb = line; s.line.width = line_w
    else:
        s.line.fill.background()
    return s


def add_text(slide, text, x, y, w, h, size=Pt(14), bold=False, color=WHITE,
             align=PP_ALIGN.LEFT, italic=False, wrap=True):
    txb = slide.shapes.add_textbox(x, y, w, h); txb.word_wrap = wrap
    tf = txb.text_frame; tf.word_wrap = wrap
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    r.font.size = size; r.font.bold = bold; r.font.italic = italic
    r.font.color.rgb = color; r.font.name = "Segoe UI"
    return txb


def add_multiline(slide, lines, x, y, w, h, spacing=Pt(6)):
    txb = slide.shapes.add_textbox(x, y, w, h); txb.word_wrap = True
    tf = txb.text_frame; tf.word_wrap = True
    for i, item in enumerate(lines):
        if isinstance(item, str):
            item = (item, GRAY, False, Pt(13))
        txt, col, bld, sz = item
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_before = spacing
        r = p.add_run(); r.text = txt
        r.font.size = sz; r.font.bold = bld; r.font.color.rgb = col
        r.font.name = "Segoe UI"
    return txb


def bg(slide, color=BG_DARK):
    add_rect(slide, 0, 0, W, H, fill=color)


def accent_line(slide, x, y, w, color=BLUE):
    ln = slide.shapes.add_shape(1, x, y, w, Inches(0.045))
    ln.fill.solid(); ln.fill.fore_color.rgb = color; ln.line.fill.background()


def tag(slide, text, x, y, color=BLUE):
    bw = Inches(len(text) * 0.125 + 0.38); bh = Inches(0.32)
    add_rect(slide, x, y, bw, bh, fill=RGBColor(color[0]//4, color[1]//4, color[2]//4),
             line=color, line_w=Pt(0.5))
    add_text(slide, text, x + Inches(0.12), y + Inches(0.035), bw, bh,
             size=Pt(9.5), bold=True, color=color, wrap=False)
    return bw


def card(slide, x, y, w, h, accent=BLUE, title=None, tsize=Pt(15)):
    add_rect(slide, x, y, w, h, fill=BG_CARD, line=GRAY_DARK, line_w=Pt(0.75))
    accent_line(slide, x, y, w, color=accent)
    if title:
        add_text(slide, title, x + Inches(0.2), y + Inches(0.16), w - Inches(0.4),
                 Inches(0.4), size=tsize, bold=True, color=WHITE)
    return y + Inches(0.6)


def dot_item(slide, text, x, y, w, dot=BLUE, size=Pt(12), color=GRAY):
    d = slide.shapes.add_shape(9, x, y + Inches(0.1), Inches(0.1), Inches(0.1))
    d.fill.solid(); d.fill.fore_color.rgb = dot; d.line.fill.background()
    add_text(slide, text, x + Inches(0.2), y, w - Inches(0.22), Inches(0.4),
             size=size, color=color)


def header(slide, label, title, sub=None, color=BLUE):
    bg(slide)
    tag(slide, label, Inches(0.55), Inches(0.5), color=color)
    add_text(slide, title, Inches(0.5), Inches(0.92), Inches(12.3), Inches(0.7),
             size=Pt(30), bold=True, color=WHITE)
    if sub:
        add_text(slide, sub, Inches(0.55), Inches(1.62), Inches(12.3), Inches(0.5),
                 size=Pt(14), color=GRAY)
    accent_line(slide, Inches(0.55), Inches(2.12), Inches(2.0), color=color)


def stat(slide, value, label, x, y, w=Inches(3.95), color=RED):
    add_text(slide, value, x, y, w, Inches(0.7), size=Pt(31), bold=True, color=color, wrap=False)
    add_text(slide, label, x, y + Inches(0.74), w, Inches(0.7), size=Pt(12), color=GRAY)


# ════════════════════════════════════════════════════════════════════════════
# 1 — TITLE
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK); bg(s)
add_rect(s, 0, 0, W, Inches(0.12), fill=CYAN)
add_text(s, "🛡", Inches(0.5), Inches(2.0), Inches(2), Inches(1.4), size=Pt(64))
add_text(s, "QALQAN AI", Inches(2.0), Inches(2.05), Inches(10), Inches(1.0),
         size=Pt(52), bold=True, color=WHITE)
add_text(s, "Қазақстандық киберқауіпсіздік + экономикалық қауіптерден қорғаныс",
         Inches(2.05), Inches(3.15), Inches(10.5), Inches(0.6), size=Pt(18), color=CYAN)
add_text(s, "Азамат үшін алғашқы тегін, нақты уақыттағы алаяқтыққа қарсы есік",
         Inches(2.05), Inches(3.7), Inches(10.5), Inches(0.5), size=Pt(13), color=GRAY)
tag(s, "ДЭР республикалық байқауы · 2026", Inches(2.05), Inches(4.35), color=BLUE)
tag(s, "Қызылорда", Inches(6.7), Inches(4.35), color=PURPLE)
add_multiline(s, [
    ("Команда:", WHITE, True, Pt(12)),
    ("Барахат Мұхтар — AITU", GRAY, False, Pt(12)),
    ("Қыдырбек Елдос — Қорқыт Ата атындағы университет", GRAY, False, Pt(12)),
], Inches(2.05), Inches(5.05), Inches(8), Inches(1.3), spacing=Pt(4))
add_text(s, "qalqan-ai-nu.vercel.app   ·   @QalqanAI_bot   ·   github.com/Kennurken/qalqan-ai",
         Inches(0.5), Inches(6.95), Inches(12.3), Inches(0.4), size=Pt(11),
         color=GRAY, align=PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════════════════════════
# 2 — ПРОБЛЕМА (экономическая угроза)
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
header(s, "МӘСЕЛЕ", "Қазақстан жыл сайын миллиардтарды жоғалтады",
       "Кибералаяқтық — нақты, өлшенетін экономикалық қауіп", color=RED)
stat(s, "16.4 млрд ₸", "Киберзиян (2025, 10 ай) — 2024-ке ×29", Inches(0.55), Inches(2.55), color=RED)
stat(s, "26 300", "Интернет-алаяқтық кейсі (+86% жылдық)", Inches(4.7), Inches(2.55), color=AMBER)
stat(s, "85 млн", "Бұғатталған алаяқ қоңырау (2025)", Inches(8.8), Inches(2.55), color=BLUE)
stat(s, "54 млрд ₸", "Қаржылық пирамидалардан зиян", Inches(0.55), Inches(4.45), color=RED)
stat(s, "31 000+", "Пирамида құрбаны (бір жылда)", Inches(4.7), Inches(4.45), color=AMBER)
stat(s, "5 схема", "= барлық алаяқтықтың 81%-ы", Inches(8.8), Inches(4.45), color=GREEN)
add_text(s, "Дереккөздер: Нацбанк Антифрод-орталығы, АРРФР/АФМ, телеком операторлары (2025)",
         Inches(0.55), Inches(6.7), Inches(12), Inches(0.4), size=Pt(11), italic=True, color=GRAY)

# ════════════════════════════════════════════════════════════════════════════
# 3 — ОЛҚЫЛЫҚ (gap)
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
header(s, "ОЛҚЫЛЫҚ", "Бүгінгі қорғаныстың бәрі — азаматқа жетпейді", color=AMBER)
y = card(s, Inches(0.55), Inches(2.5), Inches(3.95), Inches(3.4), accent=RED,
         title="🏦  BACK-END")
for i, t in enumerate(["Банк/телеком ішкі антифрод", "Транзакцияны блоктайды",
                       "Азамат сұрай алмайды", "Жабық, әр банкте бөлек"]):
    dot_item(s, t, Inches(0.75), y + Inches(0.1 + i*0.55), Inches(3.6), dot=RED)
y = card(s, Inches(4.7), Inches(2.5), Inches(3.95), Inches(3.4), accent=AMBER,
         title="⏱️  РЕАКТИВТІ")
for i, t in enumerate(["Сот шешімі бойынша тізім", "Баяу, оқиғадан кейін",
                       "Жаңа домен ілінбейді", "Алдын алу жоқ"]):
    dot_item(s, t, Inches(4.9), y + Inches(0.1 + i*0.55), Inches(3.6), dot=AMBER)
y = card(s, Inches(8.85), Inches(2.5), Inches(3.95), Inches(3.4), accent=GRAY,
         title="📄  СТАТИКАЛЫҚ")
for i, t in enumerate(["АФМ пирамида тізімі — PDF", "Эксперт үшін, азамат көрмейді",
                       "API жоқ, нақты уақыт жоқ", "Шашыраңқы дереккөздер"]):
    dot_item(s, t, Inches(9.05), y + Inches(0.1 + i*0.55), Inches(3.6), dot=GRAY)
add_rect(s, Inches(0.55), Inches(6.15), Inches(12.25), Inches(0.85), fill=BG_CARD2, line=RED)
add_text(s, "❗ Азамат сілтемені / нөмірді / SMS-ті / инвест-ұсынысты тексеретін "
         "тегін, көпшілік, нақты уақыт құралы ЖОҚ",
         Inches(0.75), Inches(6.32), Inches(12), Inches(0.5), size=Pt(14.5), bold=True, color=WHITE)

# ════════════════════════════════════════════════════════════════════════════
# 4 — ШЕШІМ
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
header(s, "ШЕШІМ", "Qalqan — азаматтың алғашқы қорғаныс есігі",
       "Қазақстанның шашыраңқы антифрод-деректерін бір қауіпсіздік тексеруіне айналдырады",
       color=CYAN)
items = [
    ("🔗 URL / сайт", "Фишинг, клон, homoglyph (kаspi.kz)"),
    ("📞 Телефон +7 7xx", "Скам-префикс, паттерн талдау"),
    ("💬 SMS / хабарлама", "Алаяқ сілтеме + мәтін талдау"),
    ("💰 Қаржылық пирамида", "АФМ реестрі + атау бойынша іздеу"),
    ("🏛️ Госзакуп-фрод", "Жалған тендер/жеткізуші белгілері"),
    ("🌐 3 тіл", "Қазақша / орысша / ағылшынша вердикт"),
]
for i, (t, d) in enumerate(items):
    col = i % 3; row = i // 3
    x = Inches(0.55 + col * 4.12); yy = Inches(2.55 + row * 1.75)
    yc = card(s, x, yy, Inches(3.9), Inches(1.55), accent=CYAN, title=t, tsize=Pt(14))
    add_text(s, d, x + Inches(0.2), yc + Inches(0.05), Inches(3.5), Inches(0.8),
             size=Pt(11.5), color=GRAY)
add_text(s, "Жеткізу: Chrome / Firefox кеңейтімі  +  Telegram бот — алаяқтық болатын жерде",
         Inches(0.55), Inches(6.45), Inches(12), Inches(0.5), size=Pt(13), bold=True, color=CYAN)

# ════════════════════════════════════════════════════════════════════════════
# 5 — ЖИВЫЕ ФИЧИ
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
header(s, "ІСКЕ ҚОСЫЛҒАН", "6 функция — бәрі live, продакшенде", color=GREEN)
feats = [
    ("Homoglyph / typosquat", "kаspi.kz (кирилл) → ҚАУІПТІ. KZ брендтерінің клондары", GREEN),
    ("АФМ-пирамида чек", "Атау бойынша: «Финико» → реестрде тіркелген", AMBER),
    ("Госзакуп-фрод", "goszakup.gov.kz — 10 ереже бойынша талдау", BLUE),
    ("Community voting", "Краудсорс: растау/жоққа шығару → авто-блок", PURPLE),
    ("Threat-feeds + KZ-фид", "URLhaus live ингест + /feed/kz (CC-BY, ашық)", CYAN),
    ("Регулятор дашборды", "KZ облыстары бойынша қауіп жылу-картасы", RED),
]
for i, (t, d, c) in enumerate(feats):
    col = i % 3; row = i // 3
    x = Inches(0.55 + col * 4.12); yy = Inches(2.5 + row * 1.85)
    yc = card(s, x, yy, Inches(3.9), Inches(1.65), accent=c, title=t, tsize=Pt(13.5))
    add_text(s, d, x + Inches(0.2), yc + Inches(0.02), Inches(3.5), Inches(0.9),
             size=Pt(11), color=GRAY)
add_text(s, "✅ 80 автотест + CI  ·  6-деңгейлі pipeline  ·  Vercel + Supabase + Upstash",
         Inches(0.55), Inches(6.55), Inches(12), Inches(0.5), size=Pt(12.5), bold=True, color=GREEN)

# ════════════════════════════════════════════════════════════════════════════
# 6 — KZ ҚАУІП КАРТАСЫ
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
header(s, "РЕГУЛЯТОР ҮШІН", "KZ нақты уақыттағы қауіп картасы",
       "Облыстар бойынша скам жылу-картасы — анонимді (url_hash/ip_hash)", color=RED)
# schematic heat tiles
regions = [("Алматы қ.", RED), ("Астана", RED), ("Шымкент", AMBER), ("Қарағанды", AMBER),
           ("Түркістан", AMBER), ("Атырау", BLUE), ("Маңғыстау", BLUE), ("Ақтөбе", BLUE),
           ("Павлодар", GRAY_DARK), ("ШҚО", GRAY_DARK), ("Қызылорда", AMBER), ("Жамбыл", BLUE)]
for i, (nm, c) in enumerate(regions):
    col = i % 4; row = i // 4
    x = Inches(0.55 + col * 3.0); yy = Inches(2.55 + row * 1.0)
    add_rect(s, x, yy, Inches(2.8), Inches(0.85), fill=c if c != GRAY_DARK else BG_CARD2,
             line=GRAY_DARK)
    add_text(s, nm, x + Inches(0.15), yy + Inches(0.22), Inches(2.5), Inches(0.4),
             size=Pt(12.5), bold=True, color=WHITE)
add_text(s, "Жүйе Vercel geo-деректерінен облысты анықтайды → регулятор үшін эконом-қауіп "
         "панелі (ДЭР тақырыбына тура)",
         Inches(0.55), Inches(6.55), Inches(12), Inches(0.5), size=Pt(12), italic=True, color=GRAY)

# ════════════════════════════════════════════════════════════════════════════
# 7 — АРХИТЕКТУРА
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
header(s, "АРХИТЕКТУРА", "6-деңгейлі анықтау pipeline", color=BLUE)
tiers = [("1", "Ақ тізім", GREEN), ("2", "Redis кэш", BLUE), ("3", "Офлайн DB", BLUE),
         ("4", "KZ Intel", PURPLE), ("5", "Сыртқы DB + домен", AMBER), ("6", "Groq / Gemini AI", CYAN)]
for i, (n, t, c) in enumerate(tiers):
    x = Inches(0.55 + i * 2.12)
    add_rect(s, x, Inches(2.6), Inches(1.9), Inches(1.4), fill=BG_CARD, line=c)
    add_text(s, n, x + Inches(0.1), Inches(2.7), Inches(1.7), Inches(0.5), size=Pt(24),
             bold=True, color=c, align=PP_ALIGN.CENTER)
    add_text(s, t, x + Inches(0.1), Inches(3.35), Inches(1.7), Inches(0.6), size=Pt(11),
             color=WHITE, align=PP_ALIGN.CENTER)
    if i < 5:
        add_text(s, "→", x + Inches(1.92), Inches(2.95), Inches(0.25), Inches(0.5),
                 size=Pt(18), color=GRAY, align=PP_ALIGN.CENTER)
yc = card(s, Inches(0.55), Inches(4.5), Inches(12.25), Inches(2.1), accent=BLUE,
          title="🧱  Технологиялық стек")
specs = ["Backend: FastAPI (Python) · Frontend: Chrome/Firefox MV3 + React",
         "Деректер: Supabase PostgreSQL · Кэш: Upstash Redis · AI: Groq llama-3.3-70b + Gemini 2.5",
         "Бот: Telegram webhook · CI/CD: GitHub Actions + Vercel · 80 автотест"]
for i, t in enumerate(specs):
    dot_item(s, t, Inches(0.8), yc + Inches(0.05 + i*0.45), Inches(11.6), dot=BLUE, size=Pt(12.5))

# ════════════════════════════════════════════════════════════════════════════
# 8 — ДАТА-MOAT
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
header(s, "НЕГЕ ҚАЙТАЛАУ ҚИЫН", "Бәсекеге төзімді артықшылықтар", color=PURPLE)
moat = [
    ("🗄️ KZ дата-moat", "Өздігінен өсетін KZ қауіп базасы + ашық CC-BY фид (/feed/kz). "
     "Жаһандық ойыншыларда жоқ, банктер жабық ұстайды."),
    ("Қазақ тілі", "Қазақша AI-вердикт пен түсіндірме. Google/Netcraft локализацияламайды."),
    ("🧩 Кеңдік", "URL + телефон + SMS + пирамида + госзакуп — бір құралда. Ешкімде жоқ комбинация."),
    ("🏛️ Ресми агрегация", "KZ-CERT + АФМ + госзакуп үстіндегі жалғыз азаматтық қабат."),
]
for i, (t, d) in enumerate(moat):
    row = i // 2; col = i % 2
    x = Inches(0.55 + col * 6.15); yy = Inches(2.5 + row * 1.95)
    yc = card(s, x, yy, Inches(5.95), Inches(1.75), accent=PURPLE, title=t, tsize=Pt(15))
    add_text(s, d, x + Inches(0.2), yc + Inches(0.02), Inches(5.6), Inches(1.0),
             size=Pt(11.5), color=GRAY)
add_text(s, "Жеткізу де moat: Telegram бот + кеңейтім — скам нақ Telegram/WhatsApp/Instagram-да тұрады",
         Inches(0.55), Inches(6.55), Inches(12), Inches(0.5), size=Pt(12), bold=True, color=PURPLE)

# ════════════════════════════════════════════════════════════════════════════
# 9 — ИМПАКТ + МАСШТАБ + B2G
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
header(s, "ИМПАКТ + МАСШТАБ", "Демодан мемлекеттік ауқымға дейін", color=GREEN)
y = card(s, Inches(0.55), Inches(2.5), Inches(3.95), Inches(3.5), accent=GREEN, title="📈 Әсер")
for i, t in enumerate(["Азамат шығынын азайту", "Регулятор үшін live аналитика",
                       "Ашық KZ қауіп фиді", "Қазақша қол жетімділік"]):
    dot_item(s, t, Inches(0.75), y + Inches(0.1 + i*0.6), Inches(3.6), dot=GREEN)
y = card(s, Inches(4.7), Inches(2.5), Inches(3.95), Inches(3.5), accent=BLUE, title="⚙️ Масштаб")
for i, t in enumerate(["10к → 1М тексеру/күн", "Тёплый кэш (VPS)", "Кезек + rate-limit",
                       "€5/ай → $50-80 @1М"]):
    dot_item(s, t, Inches(4.9), y + Inches(0.1 + i*0.6), Inches(3.6), dot=BLUE)
y = card(s, Inches(8.85), Inches(2.5), Inches(3.95), Inches(3.5), accent=AMBER, title="🏛️ B2G")
for i, t in enumerate(["Банк/МФО үшін API", "Нацбанк Антифрод-орталығын қоректендіру",
                       "KZ-CERT-ке эскалация", "eGov интеграция"]):
    dot_item(s, t, Inches(9.05), y + Inches(0.1 + i*0.6), Inches(3.6), dot=AMBER)
add_text(s, "Бәсекелесу емес — мемлекеттік антифродты ТОЛЫҚТЫРАМЫЗ (азаматтық қабат)",
         Inches(0.55), Inches(6.45), Inches(12), Inches(0.5), size=Pt(13), bold=True, color=WHITE)

# ════════════════════════════════════════════════════════════════════════════
# 10 — ДЕМО + CTA
# ════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK); bg(s)
add_rect(s, 0, 0, W, Inches(0.12), fill=CYAN)
add_text(s, "Тірі демо", Inches(0.5), Inches(0.9), Inches(12), Inches(0.8),
         size=Pt(32), bold=True, color=WHITE, align=PP_ALIGN.CENTER)
links = [("🌐 Лендинг + дашборд", "qalqan-ai-nu.vercel.app/dashboard"),
         ("🤖 Telegram бот", "@QalqanAI_bot — /check, /phone, /sms"),
         ("Ашық код", "github.com/Kennurken/qalqan-ai"),
         ("Ашық KZ фид", "qalqan-ai-nu.vercel.app/feed/kz (CC-BY)")]
for i, (t, u) in enumerate(links):
    row = i // 2; col = i % 2
    x = Inches(1.4 + col * 5.8); yy = Inches(2.1 + row * 1.25)
    yc = card(s, x, yy, Inches(5.5), Inches(1.05), accent=CYAN, title=t, tsize=Pt(15))
    add_text(s, u, x + Inches(0.2), yc + Inches(0.0), Inches(5.1), Inches(0.4),
             size=Pt(12), color=CYAN)
add_rect(s, Inches(1.4), Inches(4.8), Inches(10.5), Inches(1.3), fill=BG_CARD2, line=CYAN)
add_text(s, "«Қазақстан жыл сайын миллиардтаған теңге жоғалтады.\n"
         "Qalqan — кез келген азамат сілтемені, нөмірді, SMS-ті тексеретін алғашқы есік.»",
         Inches(1.7), Inches(5.05), Inches(10), Inches(0.9), size=Pt(15), bold=True,
         color=WHITE, align=PP_ALIGN.CENTER)
add_text(s, "Рахмет!  ·  Барахат Мұхтар + Қыдырбек Елдос  ·  ДЭР 2026",
         Inches(0.5), Inches(6.7), Inches(12.3), Inches(0.4), size=Pt(12),
         color=GRAY, align=PP_ALIGN.CENTER)

prs.save("Qalqan_AI_Defense_2026.pptx")
print("saved Qalqan_AI_Defense_2026.pptx —", len(prs.slides.__iter__.__self__._sldIdLst), "slides")
