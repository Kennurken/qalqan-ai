"""
QalqanAI — Scientific Seminar Presentation Generator
Creates a professional PowerPoint for the 2026-04-14 seminar
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# Colors
BG = RGBColor(2, 6, 23)        # #020617
SURFACE = RGBColor(15, 23, 42)  # #0f172a
CARD = RGBColor(30, 41, 59)     # #1e293b
BLUE = RGBColor(59, 130, 246)   # #3b82f6
CYAN = RGBColor(6, 182, 212)    # #06b6d4
RED = RGBColor(239, 68, 68)     # #ef4444
GREEN = RGBColor(16, 185, 129)  # #10b981
YELLOW = RGBColor(245, 158, 11) # #f59e0b
PURPLE = RGBColor(139, 92, 246) # #8b5cf6
WHITE = RGBColor(241, 245, 249) # #f1f5f9
MUTED = RGBColor(148, 163, 184) # #94a3b8
DARK_CARD = RGBColor(20, 30, 48)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

def add_bg(slide, color=BG):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_text(slide, left, top, width, height, text, size=18, color=WHITE, bold=False, align=PP_ALIGN.LEFT, font_name="Segoe UI"):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = align
    return txBox

def add_shape(slide, left, top, width, height, color=CARD):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    shape.shadow.inherit = False
    return shape

def add_card(slide, left, top, width, height, title, body, icon="", title_color=WHITE):
    add_shape(slide, left, top, width, height, CARD)
    if icon:
        add_text(slide, left + 0.2, top + 0.15, 0.6, 0.5, icon, size=28, align=PP_ALIGN.LEFT)
    add_text(slide, left + 0.2 + (0.5 if icon else 0), top + 0.15, width - 0.6, 0.4, title, size=16, bold=True, color=title_color)
    add_text(slide, left + 0.2, top + 0.65, width - 0.4, height - 0.8, body, size=12, color=MUTED)

# ============================================================
# SLIDE 1: TITLE
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)

# Shield circle
shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(5.7), Inches(0.5), Inches(1.9), Inches(1.9))
shape.fill.solid()
shape.fill.fore_color.rgb = BLUE
shape.line.fill.background()
add_text(slide, 5.9, 0.75, 1.5, 1.5, "🛡️", size=60, align=PP_ALIGN.CENTER)

add_text(slide, 2, 2.7, 9.3, 0.8, "QalqanAI", size=60, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(slide, 2, 3.5, 9.3, 0.6, "ИИ ассистент против киберугроз в интернет-среде", size=24, color=BLUE, align=PP_ALIGN.CENTER)

add_text(slide, 2, 4.5, 9.3, 0.4, "Научный семинар 2026", size=14, color=MUTED, align=PP_ALIGN.CENTER)
add_text(slide, 2, 5.0, 9.3, 0.4, "Марукоб Елдос", size=18, color=WHITE, align=PP_ALIGN.CENTER)
add_text(slide, 2, 5.5, 9.3, 0.4, "Қызылорда, 2026", size=14, color=MUTED, align=PP_ALIGN.CENTER)

# ============================================================
# SLIDE 2: PROBLEM
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text(slide, 0.5, 0.3, 3, 0.3, "ПРОБЛЕМА", size=12, bold=True, color=BLUE)
add_text(slide, 0.5, 0.6, 12, 0.8, "Киберугрозы в Казахстане: масштаб проблемы", size=36, bold=True, color=WHITE)

stats = [
    ("+340%", "Рост фишинг-атак\nв КЗ за 2024-2025", RED),
    ("₸12.8B", "Ущерб от онлайн-\nмошенничества в 2025", YELLOW),
    ("46+", "Финансовых пирамид\nработающих в КЗ", PURPLE),
    ("95%", "Фишинг-сайтов\nмладше 30 дней", RED),
]

for i, (num, label, color) in enumerate(stats):
    x = 0.5 + i * 3.1
    add_shape(slide, x, 1.8, 2.8, 2.2, CARD)
    add_text(slide, x, 2.0, 2.8, 0.8, num, size=44, bold=True, color=color, align=PP_ALIGN.CENTER)
    add_text(slide, x + 0.2, 2.9, 2.4, 1.0, label, size=14, color=MUTED, align=PP_ALIGN.CENTER)

add_text(slide, 1, 4.5, 11, 1.5,
    "Существующие решения (Google Safe Browsing, Kaspersky) не учитывают специфику\n"
    "казахстанского интернет-пространства: кириллические гомоглифы, казахоязычный\n"
    "фишинг, местные финансовые пирамиды, Telegram/WhatsApp мошенничество",
    size=16, color=MUTED, align=PP_ALIGN.CENTER)

# ============================================================
# SLIDE 3: SOLUTION — 5-Tier Pipeline
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text(slide, 0.5, 0.3, 3, 0.3, "РЕШЕНИЕ", size=12, bold=True, color=BLUE)
add_text(slide, 0.5, 0.6, 12, 0.8, "QalqanAI: 5-уровневая система защиты", size=36, bold=True, color=WHITE)

pipeline = [
    ("Cache +\nWhitelist", "<1ms", GREEN),
    ("Pyramid DB\n(46 схем)", "<5ms", YELLOW),
    ("ML Features\n(37 признаков)", "<10ms", CYAN),
    ("Threat DBs\n(4 базы)", "<500ms", PURPLE),
    ("Domain Intel\n(RDAP+SSL)", "<1s", BLUE),
    ("Groq AI\n(14,400/day)", "<3s", RED),
]

for i, (name, time, color) in enumerate(pipeline):
    x = 0.3 + i * 2.1
    add_shape(slide, x, 2.0, 1.8, 1.8, CARD)
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(2.0), Inches(1.8), Inches(0.03))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    add_text(slide, x, 2.15, 1.8, 0.35, time, size=13, bold=True, color=color, align=PP_ALIGN.CENTER)
    add_text(slide, x, 2.55, 1.8, 0.9, name, size=13, color=WHITE, align=PP_ALIGN.CENTER)

    if i < len(pipeline) - 1:
        add_text(slide, x + 1.8, 2.6, 0.3, 0.5, "→", size=20, color=BLUE, align=PP_ALIGN.CENTER)

add_text(slide, 0.5, 4.3, 12, 0.8,
    "Каскадная архитектура: каждый следующий уровень задействуется\nтолько если предыдущий не дал вердикт",
    size=16, color=MUTED, align=PP_ALIGN.CENTER)

# ============================================================
# SLIDE 4: ARCHITECTURE
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text(slide, 0.5, 0.3, 3, 0.3, "АРХИТЕКТУРА", size=12, bold=True, color=BLUE)
add_text(slide, 0.5, 0.6, 12, 0.8, "Технологический стек", size=36, bold=True, color=WHITE)

cards = [
    ("🧠", "AI Layer", "Groq Llama 3.1 8B — 14,400 req/day\nGemini 2.5 Flash — backup + Vision\nMulti-provider fallback chain"),
    ("🔍", "Threat Intelligence", "PhishTank, Google Safe Browsing\nURLhaus, OpenPhish, RDAP\nSSL Check, Pyramid DB (46 схем)"),
    ("🔬", "ML Features Engine", "37 лексических URL-признаков\nГомоглиф-детектор (Cyrillic↔Latin)\nBrand similarity (50 KZ брендов)"),
    ("🛡️", "Chrome Extension", "Manifest V3, React 19\nAuto-check каждой страницы\nБлокировка + уведомления"),
    ("📊", "XAI (Explainability)", "Factor breakdown каждого вердикта\nCounterfactual explanations\nConfidence scoring"),
    ("🌐", "API & Deploy", "FastAPI на Vercel Serverless\nRate Limiting, Input Validation\n3 языка: қаз / рус / eng"),
]

for i, (icon, title, body) in enumerate(cards):
    col = i % 3
    row = i // 3
    x = 0.5 + col * 4.15
    y = 1.6 + row * 2.7
    add_card(slide, x, y, 3.85, 2.3, title, body, icon)

# ============================================================
# SLIDE 5: ML FEATURES
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text(slide, 0.5, 0.3, 3, 0.3, "ML FEATURE ENGINEERING", size=12, bold=True, color=BLUE)
add_text(slide, 0.5, 0.6, 12, 0.8, "37 признаков URL-анализа", size=36, bold=True, color=WHITE)

add_shape(slide, 0.5, 1.6, 5.9, 4.5, CARD)
add_text(slide, 0.7, 1.7, 5.5, 0.4, "Лексические признаки", size=18, bold=True, color=WHITE)
code1 = (
    '// Извлечение за <10ms без HTTP запроса\n'
    '{\n'
    '  "url_entropy": 4.27,        // Shannon entropy\n'
    '  "is_free_tld": 1,           // .tk, .ml, .ga\n'
    '  "num_subdomains": 3,\n'
    '  "suspicious_keywords": ["login", "verify"],\n'
    '  "has_punycode": 0,\n'
    '  "digit_ratio": 0.12,\n'
    '  "brand_match": "kaspi",\n'
    '  "brand_edit_distance": 2\n'
    '}'
)
add_text(slide, 0.7, 2.2, 5.5, 3.5, code1, size=12, color=CYAN, font_name="Consolas")

add_shape(slide, 6.9, 1.6, 5.9, 4.5, CARD)
add_text(slide, 7.1, 1.7, 5.5, 0.4, "Гомоглиф-детекция (уникально для КЗ)", size=18, bold=True, color=WHITE)
code2 = (
    '// Cyrillic "а" (U+0430) vs Latin "a" (U+0061)\n'
    '// kаspi.kz → визуально = kaspi.kz\n\n'
    '{\n'
    '  "homoglyph_count": 1,\n'
    '  "has_mixed_script": 1,\n'
    '  "homoglyph_chars": ["а→a"],\n'
    '  "brand_match": "kaspi",\n'
    '  "brand_edit_distance": 0\n'
    '}\n\n'
    '// Результат: +25 к threat_score'
)
add_text(slide, 7.1, 2.2, 5.5, 3.5, code2, size=12, color=CYAN, font_name="Consolas")

# ============================================================
# SLIDE 6: XAI EXPLAINABILITY
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text(slide, 0.5, 0.3, 3, 0.3, "EXPLAINABILITY (XAI)", size=12, bold=True, color=BLUE)
add_text(slide, 0.5, 0.6, 12, 0.8, "Объяснимый ИИ: почему сайт опасен", size=36, bold=True, color=WHITE)

add_shape(slide, 1.5, 1.7, 10.3, 5, CARD)
add_text(slide, 2, 1.9, 5, 0.5, "❌ DANGEROUS — 87/100", size=24, bold=True, color=RED)

factors = [
    ("Domain age: 3 days", 70, RED, "+35"),
    ("No SSL certificate", 60, RED, "+30"),
    ("Brand similarity: kaspi (dist=2)", 50, YELLOW, "+15"),
    ("Suspicious keywords: login, verify", 30, YELLOW, "+10"),
    ("AI verdict: DANGEROUS", 14, BLUE, "+7"),
]

for i, (label, width_pct, color, impact) in enumerate(factors):
    y = 2.6 + i * 0.55
    add_text(slide, 2, y, 3.5, 0.4, label, size=13, color=WHITE)
    # Bar background
    bar_bg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.5), Inches(y + 0.05), Inches(4.5), Inches(0.3))
    bar_bg.fill.solid()
    bar_bg.fill.fore_color.rgb = RGBColor(20, 25, 40)
    bar_bg.line.fill.background()
    # Bar fill
    bar = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.5), Inches(y + 0.05), Inches(4.5 * width_pct / 100), Inches(0.3))
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()
    add_text(slide, 10.2, y, 0.8, 0.4, impact, size=13, bold=True, color=color, align=PP_ALIGN.RIGHT)

add_shape(slide, 2, 5.5, 9.3, 0.8, DARK_CARD)
add_text(slide, 2.2, 5.6, 9, 0.6,
    "Counterfactual: Сайт был бы SAFE (score <40) если:\nдомен старше 365 дней И имеет SSL сертификат",
    size=13, color=MUTED)

# ============================================================
# SLIDE 7: RESULTS
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text(slide, 0.5, 0.3, 3, 0.3, "РЕЗУЛЬТАТЫ", size=12, bold=True, color=BLUE)
add_text(slide, 0.5, 0.6, 12, 0.8, "Эффективность детекции", size=36, bold=True, color=WHITE)

# Table header
add_shape(slide, 0.5, 1.7, 12.3, 0.5, BLUE)
headers = ["Тест", "URL", "Вердикт", "Score", "Источник", "Время"]
col_widths = [1.5, 2.5, 2, 1, 2.8, 1.5]
x = 0.7
for h, w in zip(headers, col_widths):
    add_text(slide, x, 1.75, w, 0.4, h, size=13, bold=True, color=WHITE)
    x += w + 0.2

tests = [
    ("Whitelist", "kaspi.kz", "SAFE", "0", "whitelist", "<1ms", GREEN),
    ("Пирамида", "crowd1.com", "DANGEROUS", "95", "pyramid_list", "<5ms", RED),
    ("Пирамида", "finiko.com", "DANGEROUS", "95", "pyramid_list", "<5ms", RED),
    ("Gambling", "1xbet.com", "DANGEROUS", "80", "groq_ai", "~2s", RED),
    ("Free TLD", "kaspi-login.tk", "SUSPICIOUS", "42", "groq_ai", "~2s", YELLOW),
    ("Validation", "ftp://bad", "REJECTED", "—", "validation", "<1ms", BLUE),
]

for row_i, (test, url, verdict, score, source, time, v_color) in enumerate(tests):
    y = 2.35 + row_i * 0.5
    if row_i % 2 == 0:
        add_shape(slide, 0.5, y - 0.05, 12.3, 0.5, DARK_CARD)
    vals = [test, url, verdict, score, source, time]
    x = 0.7
    for j, (val, w) in enumerate(zip(vals, col_widths)):
        c = v_color if j == 2 else WHITE
        add_text(slide, x, y, w, 0.4, val, size=12, color=c, bold=(j == 2))
        x += w + 0.2

# ============================================================
# SLIDE 8: COMPARISON
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text(slide, 0.5, 0.3, 3, 0.3, "СРАВНИТЕЛЬНЫЙ АНАЛИЗ", size=12, bold=True, color=BLUE)
add_text(slide, 0.5, 0.6, 12, 0.8, "QalqanAI vs Конкуренты", size=36, bold=True, color=WHITE)

comp_headers = ["Функция", "Google SB", "Kaspersky", "Netcraft", "QalqanAI"]
comp_widths = [3, 2, 2, 2, 2.5]

add_shape(slide, 0.5, 1.7, 12.3, 0.5, BLUE)
x = 0.7
for h, w in zip(comp_headers, comp_widths):
    add_text(slide, x, 1.75, w, 0.4, h, size=13, bold=True, color=WHITE)
    x += w + 0.15

comp_data = [
    ("Фишинг детекция", "✅", "✅", "✅", "✅"),
    ("AI анализ", "❌", "❌", "Частично", "✅ Groq+Gemini"),
    ("Пирамиды / MLM", "❌", "❌", "❌", "✅ 46 схем"),
    ("Казахский язык", "❌", "❌", "❌", "✅ kk/ru/en"),
    ("Гомоглифы Cyrillic", "❌", "❌", "❌", "✅ 7 пар"),
    ("XAI объяснения", "❌", "❌", "❌", "✅ Factor+CF"),
    ("ML Features API", "❌", "❌", "❌", "✅ 37 features"),
    ("Стоимость", "Free", "$30/год", "Free", "Free"),
]

for row_i, row in enumerate(comp_data):
    y = 2.35 + row_i * 0.48
    if row_i % 2 == 0:
        add_shape(slide, 0.5, y - 0.05, 12.3, 0.48, DARK_CARD)
    x = 0.7
    for j, (val, w) in enumerate(zip(row, comp_widths)):
        c = BLUE if j == 4 and "✅" in val else WHITE
        b = j == 4
        add_text(slide, x, y, w, 0.4, val, size=12, color=c, bold=b)
        x += w + 0.15

# ============================================================
# SLIDE 9: SCIENTIFIC NOVELTY
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text(slide, 0.5, 0.3, 3, 0.3, "НАУЧНАЯ НОВИЗНА", size=12, bold=True, color=BLUE)
add_text(slide, 0.5, 0.6, 12, 0.8, "Что отличает QalqanAI", size=36, bold=True, color=WHITE)

novelty = [
    ("🇰🇿", "Первая система для КЗ", "50 казахстанских брендов в базе\n46 пирамидных схем с фокусом на КЗ\nkk/ru/en интерфейс\nПаттерны eGov/Kaspi фишинга"),
    ("🔤", "Гомоглиф-детекция", "Cyrillic↔Latin подмена символов\nа/a, е/e, о/o, с/c, р/p, у/y, х/x\nКритично для кириллических доменов\nНет аналогов в мире"),
    ("🧪", "Multi-Source Fusion", "7+ источников разведки в одном pipeline\nКаскадная архитектура\nDempster-Shafer подход\nResearch API для бенчмарков"),
    ("💡", "Explainable AI", "SHAP-подобные объяснения\nCounterfactual reasoning\nFactor breakdown + impact scores\nПубликуемо: IEEE TIFS, ACM CCS"),
]

for i, (icon, title, body) in enumerate(novelty):
    col = i % 2
    row = i // 2
    x = 0.5 + col * 6.4
    y = 1.6 + row * 2.75
    add_card(slide, x, y, 6.0, 2.4, title, body, icon)

# ============================================================
# SLIDE 10: LIVE DEMO
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text(slide, 0.5, 0.3, 3, 0.3, "ДЕМОНСТРАЦИЯ", size=12, bold=True, color=BLUE)
add_text(slide, 0.5, 0.6, 12, 0.8, "Live API", size=36, bold=True, color=WHITE)

add_shape(slide, 1, 1.7, 11.3, 4.5, CARD)
code = (
    '// POST https://qalqan-ai-nu.vercel.app/check-research\n\n'
    '{\n'
    '  "verdict": "DANGEROUS",\n'
    '  "threat_score": 87,\n'
    '  "url_features": {          // 37 ML признаков\n'
    '    "url_entropy": 4.27,\n'
    '    "is_free_tld": 1,\n'
    '    "brand_match": "kaspi",\n'
    '    "homoglyph_count": 1\n'
    '  },\n'
    '  "explanation": {            // XAI\n'
    '    "top_factors": [...],\n'
    '    "counterfactual": "SAFE if domain > 365 days",\n'
    '    "confidence": 0.87\n'
    '  },\n'
    '  "metadata": {\n'
    '    "processing_time_ms": 1243,\n'
    '    "ai_provider": "groq_ai"\n'
    '  }\n'
    '}'
)
add_text(slide, 1.3, 1.9, 10.7, 4, code, size=13, color=CYAN, font_name="Consolas")

# Badges
add_shape(slide, 3, 6.5, 3.5, 0.5, DARK_CARD)
add_text(slide, 3, 6.55, 3.5, 0.4, "🌐 qalqan-ai-nu.vercel.app", size=14, bold=True, color=BLUE, align=PP_ALIGN.CENTER)
add_shape(slide, 7, 6.5, 3.5, 0.5, DARK_CARD)
add_text(slide, 7, 6.55, 3.5, 0.4, "📦 Kennurken/qalqan-ai", size=14, bold=True, color=BLUE, align=PP_ALIGN.CENTER)

# ============================================================
# SLIDE 11: FUTURE
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text(slide, 0.5, 0.3, 3, 0.3, "ПЕРСПЕКТИВЫ", size=12, bold=True, color=BLUE)
add_text(slide, 0.5, 0.6, 12, 0.8, "Направления исследований", size=36, bold=True, color=WHITE)

future = [
    ("📊", "Dataset Paper", "Первый казахский фишинг-датасет\nNeurIPS Datasets Track"),
    ("🧠", "Fine-tuning", "XLM-RoBERTa на kk/ru/en\nACL/EMNLP"),
    ("👤", "User Study", "Влияние XAI на поведение\nCHI/SOUPS"),
    ("📈", "Temporal Eval", "Zero-day фишинг\nTemporal split"),
    ("💰", "Economic Impact", "Модель снижения ущерба\nв Казахстане"),
    ("📱", "Mobile", "Android/iOS extension\nДля мобильных"),
]

for i, (icon, title, body) in enumerate(future):
    col = i % 3
    row = i // 3
    x = 0.5 + col * 4.15
    y = 1.6 + row * 2.7
    add_card(slide, x, y, 3.85, 2.3, title, body, icon)

# ============================================================
# SLIDE 12: THANK YOU
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)

shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(5.7), Inches(0.8), Inches(1.9), Inches(1.9))
shape.fill.solid()
shape.fill.fore_color.rgb = BLUE
shape.line.fill.background()
add_text(slide, 5.9, 1.05, 1.5, 1.5, "🛡️", size=55, align=PP_ALIGN.CENTER)

add_text(slide, 2, 3.0, 9.3, 0.8, "Назарларыңызға рахмет!", size=48, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(slide, 2, 3.8, 9.3, 0.5, "Спасибо за внимание!", size=24, color=BLUE, align=PP_ALIGN.CENTER)

info = [
    ("GitHub", "Kennurken/qalqan-ai"),
    ("Live API", "qalqan-ai-nu.vercel.app"),
    ("Email", "kmarukob76@gmail.com"),
]

for i, (label, value) in enumerate(info):
    x = 2.5 + i * 3
    add_text(slide, x, 4.8, 2.5, 0.3, label, size=12, color=MUTED, align=PP_ALIGN.CENTER)
    add_text(slide, x, 5.1, 2.5, 0.4, value, size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

add_text(slide, 2, 6.2, 9.3, 0.4, "Qalqan AI v5.0 • Made in Kazakhstan 🇰🇿 • 2026", size=12, color=MUTED, align=PP_ALIGN.CENTER)


# Save
output_path = os.path.join(os.path.dirname(__file__), "QalqanAI_Presentation.pptx")
prs.save(output_path)
print(f"Presentation saved: {output_path}")
print(f"   {len(prs.slides)} slides, widescreen 16:9")
